import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.33)  # Widescreen 16:9
prs.slide_height = Inches(7.5)

# Corporate Color Palette
NAVY = RGBColor(11, 29, 58)
DARK_GRAY = RGBColor(60, 60, 60)
LIGHT_BG = RGBColor(245, 247, 250)
ACCENT_BLUE = RGBColor(0, 102, 204)
WHITE = RGBColor(255, 255, 255)

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG

# --- SLIDE 1: TITLE SLIDE (Dark Theme) ---
slide_layout = prs.slide_layouts[6] # Blank
slide1 = prs.slides.add_slide(slide_layout)
fill1 = slide1.background.fill
fill1.solid()
fill1.fore_color.rgb = NAVY

title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
tf = title_box.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "Root Cause Analysis & Case Closure"
p1.font.bold = True
p1.font.size = Pt(40)
p1.font.color.rgb = WHITE
p1.font.name = "Arial"

p2 = tf.add_paragraph()
p2.text = "Regulatory Grievance Register — Policy No: NN012603030817"
p2.font.size = Pt(22)
p2.font.color.rgb = RGBColor(180, 200, 220)
p2.font.name = "Arial"
p2.space_before = Pt(14)

p3 = tf.add_paragraph()
p3.text = "Prepared for Audit, Risk & Compliance Committee"
p3.font.size = Pt(14)
p3.font.color.rgb = RGBColor(140, 160, 180)
p3.font.name = "Arial"
p3.space_before = Pt(30)


# --- SLIDE 2: CASE OVERVIEW ---
slide2 = prs.slides.add_slide(slide_layout)
apply_background(slide2)

# Title
t_box = slide2.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.5), Inches(0.8))
tf = t_box.text_frame
p = tf.paragraphs[0]
p.text = "Case Overview & Investigation Timeline"
p.font.bold = True
p.font.size = Pt(28)
p.font.color.rgb = NAVY

# Left Table / Details Matrix
rows, cols = 6, 2
left, top, width, height = Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.5)
table_shape = slide2.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(3.5)

metadata = [
    ("Policy Number", "NN012603030817"),
    ("Insurer Vendor", "Shriram Life Insurance"),
    ("Assigned Agent", "Kasim (Mob: 8928941088)"),
    ("Grievance Category", "Operational / Intermediary Misconduct"),
    ("Regulatory TAT", "14 Days Maximum"),
    ("Resolution Status", "Resolved & Closed (Within TAT)")
]

for idx, (k, v) in enumerate(metadata):
    cell_k = table.cell(idx, 0)
    cell_k.text = k
    cell_k.text_frame.paragraphs[0].font.bold = True
    cell_k.text_frame.paragraphs[0].font.size = Pt(13)
    cell_k.text_frame.paragraphs[0].font.color.rgb = NAVY
    
    cell_v = table.cell(idx, 1)
    cell_v.text = v
    cell_v.text_frame.paragraphs[0].font.size = Pt(13)
    cell_v.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

# Right Box: Narrative Summary
r_box = slide2.shapes.add_textbox(Inches(6.75), Inches(1.8), Inches(5.8), Inches(4.5))
rtf = r_box.text_frame
rtf.word_wrap = True

rp1 = rtf.paragraphs[0]
rp1.text = "Incident Narrative"
rp1.font.bold = True
rp1.font.size = Pt(18)
rp1.font.color.rgb = ACCENT_BLUE

rp2 = rtf.add_paragraph()
rp2.text = "• Delivery Deficit:\n  Policyholder experienced a 60-day delay in receiving the promised physical copy of their insurance certificate."
rp2.font.size = Pt(14)
rp2.font.color.rgb = DARK_GRAY
rp2.space_before = Pt(12)

rp3 = rtf.add_paragraph()
rp3.text = "• Intermediary Misconduct:\n  Field Agent Kasim failed to process the request and actively avoided customer's follow-up phone inquiries."
rp3.font.size = Pt(14)
rp3.font.color.rgb = DARK_GRAY
rp3.space_before = Pt(12)

rp4 = rtf.add_paragraph()
rp4.text = "• Escalation & Resolution:\n  The customer logged a formal grievance. The operations team intervened via the insurer's central SPOC to ensure immediate physical dispatch."
rp4.font.size = Pt(14)
rp4.font.color.rgb = DARK_GRAY
rp4.space_before = Pt(12)


# --- SLIDE 3: ROOT CAUSE ANALYSIS (5-WHY) ---
slide3 = prs.slides.add_slide(slide_layout)
apply_background(slide3)

t_box3 = slide3.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.5), Inches(0.8))
tf3 = t_box3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "Root Cause Analysis: The 5-Why Investigation Framework"
p3.font.bold = True
p3.font.size = Pt(28)
p3.font.color.rgb = NAVY

whys = [
    ("Why 1 (The Trigger)", "Customer filed an official grievance because they lacked physical documentation after 60 days and the agent was entirely unreachable."),
    ("Why 2 (The Agent)", "The field agent stopped answering customer inquiries due to low local-branch performance supervision and weak intermediary field code enforcement."),
    ("Why 3 (The Logistics)", "The order bypassed the central automated fulfillment tracking sequence because it was logged as a manual field-agent hand-delivery commitment."),
    ("Why 4 (The Detection)", "Internal compliance monitors failed to flag the transaction because there was no data reconciliation tool matching active policies to physical courier dispatch IDs."),
    ("Why 5 (Root Cause)", "SYSTEMIC FAILURE: Operational dependency on manual field-agent distribution paths without an automated, core system exception trigger for delayed physical fulfillment.")
]

for idx, (label, desc) in enumerate(whys):
    y_pos = Inches(1.6 + (idx * 1.1))
    
    # Label Box
    lbl_box = slide3.shapes.add_textbox(Inches(0.75), y_pos, Inches(3.0), Inches(0.8))
    ltf = lbl_box.text_frame
    lp = ltf.paragraphs[0]
    lp.text = label
    lp.font.bold = True
    lp.font.size = Pt(14)
    if idx == 4:
        lp.font.color.rgb = RGBColor(180, 40, 40) # Highlight the root cause
    else:
        lp.font.color.rgb = ACCENT_BLUE
        
    # Description Box
    desc_box = slide3.shapes.add_textbox(Inches(3.8), y_pos, Inches(8.8), Inches(0.8))
    dtf = desc_box.text_frame
    dtf.word_wrap = True
    dp = dtf.paragraphs[0]
    dp.text = desc
    dp.font.size = Pt(14)
    dp.font.color.rgb = DARK_GRAY
    if idx == 4:
        dp.font.bold = True
        dp.font.color.rgb = NAVY


# --- SLIDE 4: CAPA PLAN ---
slide4 = prs.slides.add_slide(slide_layout)
apply_background(slide4)

t_box4 = slide4.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.5), Inches(0.8))
tf4 = t_box4.text_frame
p4 = tf4.paragraphs[0]
p4.text = "Corrective & Preventive Actions (CAPA) Log"
p4.font.bold = True
p4.font.size = Pt(28)
p4.font.color.rgb = NAVY

# Corrective Actions Box (Left)
ca_box = slide4.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(5.0))
catf = ca_box.text_frame
catf.word_wrap = True
cap1 = catf.paragraphs[0]
cap1.text = "Corrective Actions (Immediate Remediation)"
cap1.font.bold = True
cap1.font.size = Pt(18)
cap1.font.color.rgb = ACCENT_BLUE

cap2 = catf.add_paragraph()
cap2.text = "• Document Escalation & Delivery:\n  Intervened directly with Shriram Life's central SPOC to locate, print, and successfully ship the physical copy of policy NN012603030817 to the customer."
cap2.font.size = Pt(14)
cap2.font.color.rgb = DARK_GRAY
cap2.space_before = Pt(14)

cap3 = catf.add_paragraph()
cap3.text = "• Verified Grievance Closure:\n  Confirmed receipt with the policyholder, logged documentation evidence, and completed formal closure inside the 14-day regulatory TAT window."
cap3.font.size = Pt(14)
cap3.font.color.rgb = DARK_GRAY
cap3.space_before = Pt(14)

# Preventive Actions Box (Right)
pa_box = slide4.shapes.add_textbox(Inches(6.9), Inches(1.8), Inches(5.6), Inches(5.0))
patf = pa_box.text_frame
patf.word_wrap = True
pap1 = patf.paragraphs[0]
pap1.text = "Preventive Actions (Systemic Mitigations)"
pap1.font.bold = True
pap1.font.size = Pt(18)
pap1.font.color.rgb = RGBColor(0, 153, 76) # Green

pap2 = patf.add_paragraph()
pap2.text = "• Disciplinary Action Frame:\n  Logged an internal code-of-conduct infraction against Agent Kasim for formal disciplinary and field performance evaluation."
pap2.font.size = Pt(14)
pap2.font.color.rgb = DARK_GRAY
pap2.space_before = Pt(14)

pap3 = patf.add_paragraph()
pap3.text = "• Core System Exception Triggers:\n  Engineered a 20-day central systemic ageing trigger alert. The system now flags any active policy that does not contain a verified carrier tracking ID."
pap3.font.size = Pt(14)
pap3.font.color.rgb = DARK_GRAY
pap3.space_before = Pt(14)

pap4 = patf.add_paragraph()
pap4.text = "• Bypassing Intermediary Risk:\n  Migrated the shipping update infrastructure to direct-to-customer communications (SMS/Email automated alerts), eliminating manual dependencies."
pap4.font.size = Pt(14)
pap4.font.color.rgb = DARK_GRAY
pap4.space_before = Pt(14)

# Save presentation
prs.save("Grievance_RCA_Report.pptx")
print("PowerPoint file generated successfully as 'Grievance_RCA_Report.pptx'.")
